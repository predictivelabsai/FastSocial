from __future__ import annotations

import base64
import html
import io
import uuid
from datetime import timedelta

import httpx
from sqlalchemy import func, select

from fastsocial.config import settings
from fastsocial.db import session_scope
from fastsocial.models import (
    AccountMetricDaily,
    CompetitorMetricDaily,
    CompetitorProfile,
    Post,
    PostMetric,
    PostTarget,
    ReportRun,
    ReportSchedule,
    SocialAccount,
    Workspace,
    utcnow,
)
from fastsocial.storage import media_storage


def report_summary(session, workspace_id: uuid.UUID, days: int) -> dict:
    since = utcnow() - timedelta(days=days)
    latest = (
        select(
            PostMetric.post_target_id.label("target_id"),
            func.max(PostMetric.collected_at).label("latest_at"),
        )
        .group_by(PostMetric.post_target_id)
        .subquery()
    )
    rows = session.execute(
        select(PostMetric, Post, SocialAccount)
        .join(
            latest,
            (latest.c.target_id == PostMetric.post_target_id)
            & (latest.c.latest_at == PostMetric.collected_at),
        )
        .join(PostTarget, PostMetric.post_target_id == PostTarget.id)
        .join(Post, PostTarget.post_id == Post.id)
        .join(SocialAccount, PostTarget.social_account_id == SocialAccount.id)
        .where(Post.workspace_id == workspace_id, PostMetric.collected_at >= since)
    ).all()
    competitors = session.execute(
        select(CompetitorProfile, CompetitorMetricDaily)
        .join(CompetitorMetricDaily, CompetitorMetricDaily.competitor_id == CompetitorProfile.id)
        .where(
            CompetitorProfile.workspace_id == workspace_id,
            CompetitorMetricDaily.metric_date >= since.date(),
        )
        .order_by(CompetitorMetricDaily.metric_date.desc())
    ).all()
    accounts = session.execute(
        select(AccountMetricDaily, SocialAccount)
        .join(SocialAccount, AccountMetricDaily.social_account_id == SocialAccount.id)
        .where(
            SocialAccount.workspace_id == workspace_id,
            AccountMetricDaily.metric_date >= since.date(),
        )
    ).all()
    return {
        "days": days,
        "rows": rows,
        "competitors": competitors,
        "accounts": accounts,
        "totals": {
            "impressions": sum(metric.impressions for metric, _, _ in rows),
            "reach": sum(metric.reach for metric, _, _ in rows),
            "engagements": sum(
                metric.likes + metric.comments + metric.shares + metric.saves
                for metric, _, _ in rows
            ),
            "clicks": sum(metric.clicks for metric, _, _ in rows),
        },
    }


def render_report_html(workspace_name: str, report: dict) -> bytes:
    totals = report["totals"]
    top = sorted(
        report["rows"],
        key=lambda row: row[0].likes + row[0].comments + row[0].shares + row[0].saves,
        reverse=True,
    )[:10]
    post_rows = (
        "".join(
            "<tr>"
            f"<td>{html.escape(account.platform.title())}</td>"
            f"<td>{html.escape((post.content.get('text') or 'Untitled')[:180])}</td>"
            f"<td>{metric.impressions:,}</td>"
            f"<td>{metric.likes + metric.comments + metric.shares + metric.saves:,}</td>"
            f"<td>{metric.clicks:,}</td>"
            "</tr>"
            for metric, post, account in top
        )
        or '<tr><td colspan="5">No collected post metrics yet.</td></tr>'
    )
    competitor_rows = (
        "".join(
            "<tr>"
            f"<td>{html.escape(profile.display_name or profile.handle)}</td>"
            f"<td>{html.escape(profile.platform.title())}</td>"
            f"<td>{metric.followers:,}</td><td>{metric.engagement_rate:.2f}%</td>"
            "</tr>"
            for profile, metric in report["competitors"][:20]
        )
        or '<tr><td colspan="4">No competitor snapshots in this period.</td></tr>'
    )
    document = f"""<!doctype html><html><head><meta charset="utf-8"><title>{html.escape(workspace_name)} report</title>
<style>body{{font:14px system-ui;color:#183129;max-width:1000px;margin:40px auto;padding:0 24px}}h1{{font-size:28px}}.kpis{{display:grid;grid-template-columns:repeat(4,1fr);gap:12px}}.kpi{{padding:18px;background:#eef7f2;border-radius:10px}}.kpi b{{display:block;font-size:24px}}table{{width:100%;border-collapse:collapse;margin:18px 0 32px}}th,td{{padding:10px;border-bottom:1px solid #dce8e2;text-align:left}}small{{color:#657a72}}@media print{{body{{margin:0}}}}</style></head><body>
<small>FASTSOCIAL BRAND REPORT · LAST {report["days"]} DAYS</small><h1>{html.escape(workspace_name)}</h1>
<div class="kpis"><div class="kpi">Impressions<b>{totals["impressions"]:,}</b></div><div class="kpi">Reach<b>{totals["reach"]:,}</b></div><div class="kpi">Engagements<b>{totals["engagements"]:,}</b></div><div class="kpi">Clicks<b>{totals["clicks"]:,}</b></div></div>
<h2>Top content</h2><table><thead><tr><th>Network</th><th>Post</th><th>Impressions</th><th>Engagements</th><th>Clicks</th></tr></thead><tbody>{post_rows}</tbody></table>
<h2>Competitor context</h2><table><thead><tr><th>Competitor</th><th>Network</th><th>Followers</th><th>Engagement rate</th></tr></thead><tbody>{competitor_rows}</tbody></table>
<small>Generated {utcnow().strftime("%d %b %Y %H:%M UTC")} by FastSocial.</small></body></html>"""
    return document.encode()


def _top_content(report: dict, limit: int = 10) -> list:
    return sorted(
        report["rows"],
        key=lambda row: row[0].likes + row[0].comments + row[0].shares + row[0].saves,
        reverse=True,
    )[:limit]


def render_report_pdf(workspace_name: str, report: dict) -> bytes:
    """Render a portable, branded report without browser-side print conversion."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    output = io.BytesIO()
    document = SimpleDocTemplate(
        output,
        pagesize=landscape(A4),
        rightMargin=14 * mm,
        leftMargin=14 * mm,
        topMargin=12 * mm,
        bottomMargin=12 * mm,
        title=f"{workspace_name} social performance report",
    )
    styles = getSampleStyleSheet()
    totals = report["totals"]
    story = [
        Paragraph("FASTSOCIAL BRAND REPORT", styles["Normal"]),
        Paragraph(html.escape(workspace_name), styles["Title"]),
        Paragraph(f"Performance for the last {report['days']} days", styles["Normal"]),
        Spacer(1, 7 * mm),
    ]
    kpis = Table(
        [
            ["Impressions", "Reach", "Engagements", "Clicks"],
            [
                f"{totals['impressions']:,}",
                f"{totals['reach']:,}",
                f"{totals['engagements']:,}",
                f"{totals['clicks']:,}",
            ],
        ],
        colWidths=[62 * mm] * 4,
    )
    kpis.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#eef7f2")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#657a72")),
                ("TEXTCOLOR", (0, 1), (-1, 1), colors.HexColor("#183129")),
                ("FONTSIZE", (0, 1), (-1, 1), 18),
                ("FONTNAME", (0, 1), (-1, 1), "Helvetica-Bold"),
                ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#dce8e2")),
                ("INNERGRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#dce8e2")),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("TOPPADDING", (0, 0), (-1, -1), 7),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
            ]
        )
    )
    story.extend([kpis, Spacer(1, 7 * mm), Paragraph("Top content", styles["Heading2"])])
    content_rows = [["Network", "Post", "Impressions", "Engagements", "Clicks"]]
    for metric, post, account in _top_content(report):
        text = (post.content.get("text") or "Untitled")[:115]
        content_rows.append(
            [
                account.platform.title(),
                Paragraph(html.escape(text), styles["BodyText"]),
                f"{metric.impressions:,}",
                f"{metric.likes + metric.comments + metric.shares + metric.saves:,}",
                f"{metric.clicks:,}",
            ]
        )
    if len(content_rows) == 1:
        content_rows.append(["—", "No collected post metrics yet.", "—", "—", "—"])
    table = Table(
        content_rows, colWidths=[28 * mm, 132 * mm, 31 * mm, 31 * mm, 27 * mm], repeatRows=1
    )
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#183129")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#dce8e2")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ]
        )
    )
    story.extend(
        [
            table,
            Spacer(1, 6 * mm),
            Paragraph(
                f"Generated {utcnow().strftime('%d %b %Y %H:%M UTC')} by FastSocial.",
                styles["Normal"],
            ),
        ]
    )
    document.build(story)
    return output.getvalue()


def render_report_pptx(workspace_name: str, report: dict) -> bytes:
    """Create an editable client deck with native PowerPoint elements."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = f"{workspace_name}\nSocial performance"
    title_slide.placeholders[
        1
    ].text = (
        f"Last {report['days']} days · Generated by FastSocial · {utcnow().strftime('%d %b %Y')}"
    )

    kpi_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    kpi_slide.shapes.title.text = "Performance at a glance"
    totals = report["totals"]
    for index, (label, value) in enumerate(
        [
            ("Impressions", totals["impressions"]),
            ("Reach", totals["reach"]),
            ("Engagements", totals["engagements"]),
            ("Clicks", totals["clicks"]),
        ]
    ):
        left = Inches(0.65 + index * 3.15)
        shape = kpi_slide.shapes.add_textbox(left, Inches(2.25), Inches(2.7), Inches(1.5))
        frame = shape.text_frame
        frame.text = label
        frame.paragraphs[0].font.size = Pt(15)
        frame.paragraphs[0].font.color.rgb = RGBColor(101, 122, 114)
        number = frame.add_paragraph()
        number.text = f"{value:,}"
        number.font.size = Pt(30)
        number.font.bold = True
        number.font.color.rgb = RGBColor(24, 49, 41)

    content_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    content_slide.shapes.title.text = "Top content"
    top = _top_content(report, 7)
    shape = content_slide.shapes.add_table(
        max(2, len(top) + 1), 5, Inches(0.45), Inches(1.45), Inches(12.4), Inches(4.9)
    )
    table = shape.table
    headers = ["Network", "Post", "Impressions", "Engagements", "Clicks"]
    widths = [1.35, 6.25, 1.65, 1.75, 1.4]
    for column, (header, width) in enumerate(zip(headers, widths, strict=True)):
        table.columns[column].width = Inches(width)
        table.cell(0, column).text = header
    if top:
        for row_index, (metric, post, account) in enumerate(top, 1):
            values = [
                account.platform.title(),
                (post.content.get("text") or "Untitled")[:130],
                f"{metric.impressions:,}",
                f"{metric.likes + metric.comments + metric.shares + metric.saves:,}",
                f"{metric.clicks:,}",
            ]
            for column, value in enumerate(values):
                table.cell(row_index, column).text = value
    else:
        table.cell(1, 0).text = "No data"
        table.cell(1, 1).text = "Collect post metrics to populate this slide."
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(24, 49, 41)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True

    competitor_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    competitor_slide.shapes.title.text = "Competitor context"
    body = competitor_slide.placeholders[1].text_frame
    body.clear()
    competitors = report["competitors"][:8]
    if not competitors:
        body.paragraphs[0].text = "No competitor snapshots are available for this period."
    for index, (profile, metric) in enumerate(competitors):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = (
            f"{profile.display_name or profile.handle} · {profile.platform.title()} · "
            f"{metric.followers:,} followers · {metric.engagement_rate:.2f}% engagement"
        )
        paragraph.level = 0
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def report_json(workspace_name: str, report: dict) -> dict:
    """A stable flat-ish feed for BI tools, MCP consumers, and custom dashboards."""
    return {
        "schema_version": "1.0",
        "workspace": workspace_name,
        "generated_at": utcnow().isoformat(),
        "period_days": report["days"],
        "summary": dict(report["totals"]),
        "posts": [
            {
                "post_id": str(post.id),
                "platform": account.platform,
                "account": account.display_name or account.username,
                "text": post.content.get("text") or "",
                "impressions": metric.impressions,
                "reach": metric.reach,
                "likes": metric.likes,
                "comments": metric.comments,
                "shares": metric.shares,
                "saves": metric.saves,
                "clicks": metric.clicks,
                "collected_at": metric.collected_at.isoformat(),
            }
            for metric, post, account in report["rows"]
        ],
        "competitors": [
            {
                "competitor_id": str(profile.id),
                "name": profile.display_name or profile.handle,
                "platform": profile.platform,
                "date": metric.metric_date.isoformat(),
                "followers": metric.followers,
                "engagement_rate": metric.engagement_rate,
            }
            for profile, metric in report["competitors"]
        ],
        "accounts": [
            {
                "account_id": str(account.id),
                "platform": account.platform,
                "date": metric.metric_date.isoformat(),
                "followers": metric.followers,
                "impressions": metric.impressions,
                "engagement": metric.engagement,
            }
            for metric, account in report["accounts"]
        ],
    }


def render_scheduled_report(workspace_name: str, report: dict, output_format: str):
    formats = {
        "html": (render_report_html, "html", "text/html; charset=utf-8"),
        "pdf": (render_report_pdf, "pdf", "application/pdf"),
        "pptx": (
            render_report_pptx,
            "pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
    }
    renderer, extension, content_type = formats.get(output_format, formats["html"])
    return renderer(workspace_name, report), extension, content_type


async def execute_report_schedule(schedule_id: uuid.UUID) -> ReportRun:
    cfg = settings()
    with session_scope() as session:
        schedule = session.get(ReportSchedule, schedule_id)
        if not schedule:
            raise ValueError("Report schedule not found")
        workspace = session.get(Workspace, schedule.workspace_id)
        run = ReportRun(schedule_id=schedule.id, recipients=list(schedule.recipients))
        session.add(run)
        session.flush()
        report = report_summary(session, workspace.id, schedule.report_days)
        body, extension, content_type = render_scheduled_report(
            workspace.name, report, schedule.output_format
        )
        key = f"fastsocial/{workspace.id}/reports/{run.id}.{extension}"
        try:
            media_storage().put(key, body, content_type)
            run.storage_key = key
            run.status = "generated"
            if cfg.postmark_server_token and schedule.recipients:
                async with httpx.AsyncClient(timeout=30) as client:
                    response = await client.post(
                        "https://api.postmarkapp.com/email",
                        headers={
                            "X-Postmark-Server-Token": cfg.postmark_server_token,
                            "Content-Type": "application/json",
                        },
                        json={
                            "From": cfg.report_from_email,
                            "To": ",".join(schedule.recipients),
                            "Subject": f"{workspace.name} · {schedule.name}",
                            "HtmlBody": (
                                body.decode()
                                if extension == "html"
                                else "<p>Your FastSocial brand report is attached.</p>"
                            ),
                            "Attachments": (
                                []
                                if extension == "html"
                                else [
                                    {
                                        "Name": f"{workspace.slug}-report.{extension}",
                                        "Content": base64.b64encode(body).decode(),
                                        "ContentType": content_type,
                                    }
                                ]
                            ),
                            "MessageStream": "outbound",
                        },
                    )
                response.raise_for_status()
                run.status = "delivered"
            run.completed_at = utcnow()
            schedule.last_run_at = run.completed_at
            schedule.next_run_at = run.completed_at + timedelta(
                days=7 if schedule.frequency == "weekly" else 30
            )
        except Exception as exc:  # noqa: BLE001
            run.status = "failed"
            run.error_message = str(exc)
            run.completed_at = utcnow()
        return run


async def run_due_reports(limit: int = 10) -> int:
    with session_scope() as session:
        query = (
            select(ReportSchedule.id)
            .where(
                ReportSchedule.active.is_(True),
                ReportSchedule.next_run_at.is_not(None),
                ReportSchedule.next_run_at <= utcnow(),
            )
            .order_by(ReportSchedule.next_run_at)
            .limit(limit)
        )
        ids = list(session.scalars(query))
    for schedule_id in ids:
        await execute_report_schedule(schedule_id)
    return len(ids)
